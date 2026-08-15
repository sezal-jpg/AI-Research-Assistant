from pathlib import Path
from pptx import Presentation
from langchain_core.documents import Document
from app.core.logger import logger

class PptxLoader:
    def load(self,file_path:Path):
        logger.info(f'Loading PPTX: {file_path.name}')
        presentation=Presentation(file_path)
        docs=[]
        for slide_number,slide in enumerate(presentation.slides,start=1):
            slide_parts=[]
            for shape in slide.shapes:
                if hasattr(shape,'text'):
                    text=shape.text.strip()
                    if text:
                        slide_parts.append(text)
                        
                if shape.has_table:
                    for row in shape.table.rows:
                        cells=[cell.text.strip() for cell in row.cells]
                        slide_parts.append("|".join(cells))
                        
            slide_text= '\n'.join(slide_parts)    
                   
            if slide_text.strip():
                docs.append(Document(page_content=slide_text,metadata={
                    'source_file':file_path.name,
                    'source_type':'pptx',
                    'slide':slide_number,
                }))       
                
        logger.info(f'Loaded{len(docs)} PPTX slides')
        return docs
    
pptx_loader=PptxLoader()                 
                                